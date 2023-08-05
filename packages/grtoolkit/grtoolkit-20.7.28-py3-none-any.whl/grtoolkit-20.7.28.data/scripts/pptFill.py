# python = "C:\Users\gabrielr\AppData\Local\Programs\Python\Python38-32\python.exe"
# pip = "C:\Users\gabrielr\AppData\Local\Programs\Python\Python38-32\python.exe\Scripts\pip.exe"

import sys, os
from pptx import Presentation
# from os.path import dirname
from grtoolkit.Storage import File
from grtoolkit.PowerPoint import duplicate_slide

# parent = dirname(sys.argv[0])
parent = os.getcwd()
QUOTES = f'{parent}\\Quotes.txt'
OLD_PPT = f'{parent}\\QUOTE_TEMPLATE.pptm'
NEW_PPT = f'{parent}\\QUOTE_NEW.pptm'

prs = Presentation(OLD_PPT)

QUOTESContent = File(QUOTES).read()
QUOTESList = QUOTESContent.split("\n\n")

#DUPLICATE PPT SLIDE FOR EVERY 
# for i in range (1,len(QUOTESList)):
#     duplicate_slide(prs,0)

# i = 0
# for quote in QUOTESList:
#     if i>0:
#         duplicate_slide(prs,0)
#     i=i+1

#DUPLICATE SLIDE 0 FROM TEMPLATE
duplicate_slide(prs,0,len(QUOTESList)-1)
#CHANGE TEXT IN ALL SLIDES AS PER NOTEPAD FILE
i=0
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            shape.text_frame.clear()
            p=shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = QUOTESList[i]
            font=run.font
            font.name='Gabriola'
    i=i+1
prs.save(NEW_PPT)