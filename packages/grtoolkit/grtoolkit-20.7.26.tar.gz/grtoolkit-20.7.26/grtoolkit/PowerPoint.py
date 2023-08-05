# python = "C:\Users\gabrielr\AppData\Local\Programs\Python\Python38-32\python.exe"
# pip = "C:\Users\gabrielr\AppData\Local\Programs\Python\Python38-32\python.exe\Scripts\pip.exe"

from pptx import Presentation
import sys, copy, six
from os.path import dirname
from grtoolkit.Storage import File

# QUOTES = f'{parent}\\Quotes.txt'
# OLD_PPT = f'{parent}\\4HR.pptm'
# NEW_PPT = f'{parent}\\new-file-name.pptm'

# parent = dirname(sys.argv[0])
# prs = Presentation(OLD_PPT)
# prs.save(NEW_PPT)

def get_blank_slide_layout(pres):
    layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return pres.slide_layouts[blank_layout_id]

def duplicate_slide(pres, index, no_of_copies=1):
    """Duplicate the slide with the given index in pres.
        Adds slide to the end of the presentation"""

    source = pres.slides[index]
    for _ in range(0,no_of_copies):

        blank_slide_layout = get_blank_slide_layout(pres)
        dest = pres.slides.add_slide(blank_slide_layout)

        for shp in source.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

        for _, value in six.iteritems(source.part.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if not "notesSlide" in value.reltype:
                dest.part.rels.add_relationship(value.reltype, value._target, value.rId)

    return dest


# QUOTESContent = File(QUOTES).read()
# QUOTESList = QUOTESContent.split("\n\n")
# i = 0
# for quote in QUOTESList:
#     if i>0:
#         duplicate_slide(prs,0)
#     i=i+1

# #CHANGE TEXT IN ALL SLIDES
# i=0
# for slide in prs.slides:
#     for shape in slide.shapes:
#         if hasattr(shape, "text"):
#             shape.text_frame.clear()
#             p=shape.text_frame.paragraphs[0]
#             run = p.add_run()
#             run.text = QUOTESList[i]
#             font=run.font
#             font.name='Gabriola'
#     i=i+1
# prs.save(NEW_PPT)